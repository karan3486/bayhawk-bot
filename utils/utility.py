import smtplib
from email.mime.text import MIMEText
from email.mime.multipart import MIMEMultipart
from database import db
import json
from openai import OpenAI
from flask import Flask, render_template, request, jsonify,send_file,render_template_string
from config import Config
import os
from fpdf import FPDF 
import pdfplumber
from docx import Document
from pptx.util import Pt
from pptx import Presentation
import requests
from bs4 import BeautifulSoup

cred = Config()
client = OpenAI()
model = 'gpt-3.5-turbo'



def url_to_text(url):
    try:
        # Fetch the content of the URL
        response = requests.get(url)
        response.raise_for_status()  # Raise an error for bad status codes
        
        # Parse the HTML content
        soup = BeautifulSoup(response.text, 'html.parser')
        
        # Extract the text from the parsed HTML
        text = soup.get_text(separator="\n", strip=True)
        return text
    except requests.exceptions.RequestException as e:
        return f"An error occurred: {e}"
def startMailSending(email=None):
    conversation_datas = db.get_all_conversations()
    for conversation_data in conversation_datas:
        visitor_email = conversation_data["email"]
        visitor_name = conversation_data["name"]
        conversation = conversation_data["conversation"]
        if email:
            if email.lower()==visitor_email.lower():
                sendMailLead(conversation,visitor_name,visitor_email)
                return
        else:
            sendMailLead(conversation,visitor_name,visitor_email)

def sendMailLead(conversation,name,email):
    

    response_email = client.chat.completions.create(
    model=model,
    messages=[
            {"role": "system", "content": "You are an assistant admission staff for the SFBU(San francisco bay university) who take care of incoming new students and manage and help."},
                {
                    "role": "user",
                    "content": f''' Asuming that you provide university admission support for SFBU(San francisco bay university).
                    Create an email to be sent to the Admission Head team based on the customer interaction chat with the Conversation chat ==>: {conversation},
                    Summarize the chat conversation and create an email.
                    Details of visitor is name=>{name}, email=>{email}. Mention the details of visitor also in the email. Give only the email body content with greetings and Best regards Bayhawk Assistant Bot, and do not give email subject
                    '''
                   
                }
            ],
    )

    email_content = response_email.choices[0].message.content

    response_subject = client.chat.completions.create(
    model=model,
    messages=[
            {"role": "system", "content": "You are a helpful assistant."},
                {
                    "role": "user",
                    "content": f"Give me email subject only for this email=>{email_content}. Donot give tag-> 'Email Subject:', just give directly the subject content only."
                }
            ],
    max_tokens=100
    )
    email_subject = response_subject.choices[0].message.content
    send_email(email_content,cred.ADMISSION_TEAM,email_subject,)

def Langtranslate(text,language):
    response_subject = client.chat.completions.create(
    model=model,
    messages=[
            {"role": "system", "content": "You are a helpful assistant.An Expert in lnaguage translation"},
                {
                    "role": "user",
                    "content": f"Translate this content in language=>{language}. If the content is in English then dont translate give as it is. Content=>{text}"
                }
            ]
    )
    return response_subject.choices[0].message.content

def set_font_size(placeholder, font_size):
    """Sets the font size for all text in a placeholder."""
    if placeholder and placeholder.text_frame:
        for paragraph in placeholder.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(font_size)


def savePPTFile(slides,filename):
            results_path = os.path.join(cred.RESULTS_FOLDER, filename)
            presentation = Presentation()
            for slide_content in slides:
                slide = presentation.slides.add_slide(presentation.slide_layouts[1])  # Title and Content layout
                if ":" in slide_content:  # Assuming "Title: Content" structure
                    title, content = slide_content.split(":", 1)
                    slide.shapes.title.text = title.strip()
                    set_font_size(slide.shapes.title, 30)  # Set font size for the title
                    content_placeholder = slide.placeholders[1]
                    content_placeholder.text = content.strip()
                    set_font_size(content_placeholder, 20)  # Set font size for the content
                else:
                    slide.shapes.title.text = "Details"
                    # set_font_size(slide.shapes.title, 20)  # Set font size for the title
                    content_placeholder = slide.placeholders[1]
                    content_placeholder.text = slide_content.strip()
                    set_font_size(content_placeholder, 20)  # Set font size for the content
            presentation.save(results_path)

import smtplib
from email.mime.text import MIMEText
from email.mime.multipart import MIMEMultipart

def send_email(body,email,subject,isHtml=False):
    # Create the email message
    msg = MIMEMultipart()
    msg["To"] = email
    msg["Subject"] = subject
    if isHtml:
        msg.attach(MIMEText(body, "html"))
    else:
        msg.attach(MIMEText(body, "plain"))


    try:
        # Connect to the SMTP server
        with smtplib.SMTP("smtp.gmail.com", 587) as server:
            server.starttls()  # Secure the connection
            server.login(cred.SMTP_USERNAME , cred.SMTP_PASSWORD)  # Log in to the email account
            server.send_message(msg)  # Send the email
            print(f"Email sent to {email} successfully!")
    except Exception as e:
        print(f"Failed to send email to {email}: {e}")
# startMailSending()

def allowed_file(filename):
    return ['.' in filename and filename.rsplit('.', 1)[1].lower() in cred.ALLOWED_EXTENSIONS]

def extract_text_from_file(file_path):
    ext = file_path.rsplit('.', 1)[1].lower()
    if ext == 'pdf':
        with pdfplumber.open(file_path) as pdf:
            text = ''.join([page.extract_text() for page in pdf.pages])
        return text
    elif ext == 'docx':
        doc = Document(file_path)
        text = ' '.join([para.text for para in doc.paragraphs])
        return text
    elif ext == 'txt':
        with open(file_path, 'r') as file:
            return file.read()
    return None

def Question_mcqs_generator(input_text, num_questions,level):
    prompt = f"""
    You are an AI assistant, an expert MCQ maker helping the user generate multiple-choice questions (MCQs) based on the following text:
    '{input_text}'
    Please generate {num_questions} MCQs from the text.Generate the the Question of level: '{level}'
    Make sure the questions are not repeated and check all the questions to be conforming the text as well. Each question should have:
    - A clear question
    - Four answer options (labeled A, B, C, D)
    - The correct answer clearly indicated
    Format:
    ## MCQ
    Question: [question]
    A) [option A]
    B) [option B]
    C) [option C]
    D) [option D]
    Correct Answer: [correct option]
    """
    response = client.chat.completions.create(
        model="gpt-4",
        messages=[{"role": "system", "content": "You are an expert in creating MCQs."},
                  {"role": "user", "content": prompt}]
    )
    return response.choices[0].message.content.strip()

def save_mcqs_to_file(mcqs, filename):
    results_path = os.path.join(cred.RESULTS_FOLDER, filename)
    with open(results_path, 'w') as f:
        f.write(mcqs)
    return results_path

def create_pdf(mcqs, filename):
    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("Arial", size=12)

    for mcq in mcqs.split("## MCQ"):
        if mcq.strip():
            pdf.multi_cell(0, 10, mcq.strip())
            pdf.ln(5)  # Add a line break

    pdf_path = os.path.join(cred.RESULTS_FOLDER, filename)
    pdf.output(pdf_path)
    return pdf_path

def getPPTContent(context,no_slides):
    prompt = f'''You are an expert at generating professional and well-structured PowerPoint presentations. 
    Your task is to produce detailed slide content based on the provided input context. Create total number of slides:{no_slides}
    
    Follow these rules:

1. **Structure:** 
   - Each slide should have a title and bullet points.
   - Use clear, concise, and professional language.

2. **Slide Types:** 
   - Include slides for an introduction, main points, detailed explanations, and a conclusion.

3. **Formatting:**
   - Each slide should have relavant topic caption.
   - Slide titles should summarize the content of the slide.
   - Use 3 to 4 bullet points per slide.
   - Each bullet point should contain no more than two sentences.

4. **Content Style:**
   - Use formal, professional language.
   - Avoid jargon unless necessary and explain it if used.
   - Maintain consistency across slides.

5. **Output Format:**
   - Provide the content for each slide as:
     "Slide Topic [Slide Title]: \n- [Bullet Point 1]\n- [Bullet Point 2]\n...".
   - Separate each slide with a blank line.

6. **Example:**
   For input "Overview of AI in Healthcare," generate:
   Slide 1 - Introduction to AI in Healthcare:
   - Artificial Intelligence (AI) is transforming healthcare.
   - AI applications range from diagnostics to personalized treatment.

   Slide 2 - Benefits of AI in Healthcare:
   - Improves diagnostic accuracy through advanced algorithms.
   - Enhances patient care with personalized treatment plans.
   - Reduces operational costs in healthcare systems.

   Slide 3 - Conclusion:
   - Reduces operational costs in healthcare systems.
   - And increase the chance of accuracy.

7. Always begin by acknowledging the topic provided in the input.
6. Always end with the conclusion slide.

\n
Context:{context}

'''
    response = client.chat.completions.create(
        model=model,
        messages=[{"role": "system", "content": "You are an expert in creating Power Point Slides."},
                  {"role": "user", "content": prompt}]
    )
    return response.choices[0].message.content.strip()

def getSyllabus(course,weekly_session,hr_session,duration):
    system_prompt = f'''You are an expert academic advisor specializing in designing detailed and well-structured course syllabi for university-level courses. Your role is to assist faculty members in creating a comprehensive syllabus based on the context they provide. Follow these guidelines:

1. **Input Context**:
   - The course or subject name.
   - The number of class sessions per week.
   - The length of each class session (in hours).
   - The total duration of the semester (in months).

2. **Output Requirements**:
   - Create a syllabus that fits within the provided timeframe and ensures the course is fully covered by the end of the semester.
   - Divide the syllabus into weekly topics, ensuring logical progression.
   - Allocate time for key sections such as:
     - Introduction and orientation in the first week.
     - Midterm review and assessment at the semester's midpoint divide the months by 2 and make midtern in the mid or months.
     - Final review and assessment in the last week.

3. **Topic Selection**:
   - Include the key topics reprompte, discussion, practical)
   - For example:
     ```
     ##section
     Course: Introduction to Artificial Intelligence
     Duration: 4 months (3 sessions per week, 2 hours per session)
     
     ##section
     Week 1:
       - Topic: Introduction to AI
       - Activities: Overview of AI, History, Applications, Discussion
       
     ##section
     Week 2:
       - Topic: Basics of Machine Learning
       - Activities: Supervised Learning, Hands-on Example, Q&A
     ```

6. **Language and Style**:
   - Use clear, formal, and professional language.
   - Ensure that the output is easy to understand and actionable for the faculty.

7. **Constraints**:
   - The syllabus must be acweekly_sessionhievable within the provided time frame.
   - Avoid including overly advanced topics that are not appropriate for the course level.

8. **Acknowledgment of Context**:
   - Begin by summarizing the provided input (course name, weekly sessions, hours, and semester duration) before generating the syllabus.

'''
    
    context = f"""
    Course: {course}
    Weekly sessions: {weekly_session}
    Hours per session: {hr_session}
    Semester duration: {duration} months
        """
    response = client.chat.completions.create(
        model=model,
        messages=[{"role": "system", "content": "You are an expert in creating Power Point Slides."},
                  {"role": "user", "content": f"{system_prompt}\n\nInput Context: {context}\n\nGenerate a syllabus:"}]
    )
    return response.choices[0].message.content.strip()

def handle_moderation_response(moderation_output):
    # Extract moderation results
    result = moderation_output["results"][0]
    flagged = result["flagged"]
    categories = result["categories"]
    category_scores = result["category_scores"]
    
    if flagged:
        # Construct response for flagged content
        response = {
            "status": "error",
            "message": "Your input has been flagged for violating content policies.",
            "flagged_categories": [
                category for category, flagged in categories.items() if flagged
            ],
            "category_scores": {
                category: score for category, score in category_scores.items() if score > 0.5
            }
        }
    else:
        # Construct response for non-flagged content
        response = {
            "status": "success",
            "message": "Your input passed moderation checks.",
        }
    
    return json.dumps(response, indent=4)
